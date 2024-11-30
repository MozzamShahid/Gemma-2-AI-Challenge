import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

def process_document(file_path):
    file_extension = file_path.split('.')[-1].lower()

    if file_extension == 'pdf':
        return process_pdf(file_path)
    elif file_extension == 'docx':
        return process_docx(file_path)
    elif file_extension == 'pptx':
        return process_pptx(file_path)
    elif file_extension in ['png', 'jpg', 'jpeg']:
        return process_image(file_path)
    else:
        return "Unsupported file format"

def process_pdf(file_path):
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

def process_docx(file_path):
    doc = Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def process_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
    return text

def process_image(file_path):
    return pytesseract.image_to_string(Image.open(file_path))